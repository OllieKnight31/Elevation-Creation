#!/usr/bin/env python3
"""
Build Property Labs Webinar PPTX from HTML slide data.
170 slides, dark green-black background (#0a0f0d), green accent (#3ECFA0).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours ──────────────────────────────────────────────────────────────────
BG       = RGBColor(10, 15, 13)       # #0a0f0d
WHITE    = RGBColor(255, 255, 255)
WHITE_M  = RGBColor(179, 179, 179)    # ~70% white
GREEN    = RGBColor(62, 207, 160)     # #3ECFA0
GREEN2   = RGBColor(43, 184, 138)     # #2BB88A

# ── Slide dimensions: 16:9 widescreen ────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def add_textbox(slide, text, left, top, width, height,
                font_size=32, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multi_line(slide, lines, left, top, width, height,
                   base_size=32, align=PP_ALIGN.LEFT):
    """
    lines: list of (text, font_size, bold, italic, color) tuples
    Each tuple becomes one paragraph.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, fsize, fbold, fitalic, fcolor) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name   = "Calibri"
        run.font.size   = Pt(fsize)
        run.font.bold   = fbold
        run.font.italic = fitalic
        run.font.color.rgb = fcolor
    return txBox


def center_box(slide, lines, v_offset_pct=0.0):
    """Centre-stacked slide. lines as above. v_offset_pct shifts centre up/down."""
    margin = Inches(1.2)
    left = margin
    top = Inches(0.6 + 7.5 * v_offset_pct)
    width = W - 2 * margin
    height = H - Inches(1.2)
    return add_multi_line(slide, lines, left, top, width, height,
                          align=PP_ALIGN.CENTER)


def left_box(slide, lines, v_offset_pct=0.0):
    """Left-aligned slide."""
    left = Inches(1.4)
    top = Inches(1.0 + 7.5 * v_offset_pct)
    width = W - Inches(2.4)
    height = H - Inches(1.5)
    return add_multi_line(slide, lines, left, top, width, height,
                          align=PP_ALIGN.LEFT)


def label_line(text, size=13):
    return (text.upper(), size, True, False, GREEN)

def label_line_gold(text, size=13):
    return (text.upper(), size, True, False, GREEN)

def white_line(text, size=36):
    return (text, size, False, False, WHITE)

def muted_line(text, size=36):
    return (text, size, False, False, WHITE_M)

def green_line(text, size=36):
    return (text, size, True, False, GREEN)

def italic_white(text, size=36):
    return (text, size, False, True, WHITE)

def italic_green(text, size=36):
    return (text, size, False, True, GREEN)

def spacer(size=12):
    return ("", size, False, False, WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def section_divider(label_text, title_text, sub_text):
    s = new_slide()
    center_box(s, [
        label_line(label_text, 13),
        spacer(10),
        (title_text, 80, True, False, GREEN),
        spacer(8),
        (sub_text, 20, False, False, WHITE_M),
    ])
    return s


def big_number_slide(number, label):
    s = new_slide()
    center_box(s, [
        (number, 100, True, False, GREEN),
        spacer(8),
        (label, 22, False, False, WHITE_M),
    ])
    return s


def case_study_header(name, sub):
    s = new_slide()
    center_box(s, [
        label_line("Case Study", 13),
        spacer(8),
        (name, 64, True, False, GREEN),
        spacer(8),
        (sub, 28, False, False, WHITE),
    ])
    return s


def case_study_detail(tag, lines):
    """lines: list of (text, color)"""
    s = new_slide()
    built = [label_line(tag, 13), spacer(8)]
    for (t, c) in lines:
        built.append((t, 30, True if c == WHITE else False, False, c))
    left_box(s, built)
    return s


def misconception(number, claim, verdict, proof=None):
    s = new_slide()
    lines = [
        (number, 80, True, False, RGBColor(30, 60, 45)),   # very dim green
        spacer(4),
        (claim, 32, False, False, WHITE_M),
        spacer(4),
        (verdict, 48, True, False, GREEN),
    ]
    if proof:
        lines.append(spacer(6))
        lines.append((proof, 18, False, False, WHITE_M))
    left_box(s, lines)
    return s


def objection_slide(question, answer):
    s = new_slide()
    left_box(s, [
        (question, 26, False, False, WHITE_M),
        spacer(8),
        ("─" * 20, 12, False, False, GREEN),
        spacer(8),
        (answer, 32, True, False, WHITE),
    ])
    return s


# ─────────────────────────────────────────────────────────────────────────────
# BUILD ALL 170 SLIDES
# ─────────────────────────────────────────────────────────────────────────────

# SLIDE 1 — Brand intro
s = new_slide()
center_box(s, [
    label_line("Property Labs Presents", 13),
    spacer(8),
    ("The 3C System:", 64, True, False, GREEN),
    spacer(6),
    ("How Complete Beginners Are Launching\nProfitable Airbnbs in 90 Days", 36, True, False, WHITE),
    spacer(20),
    ("PROPERTY LABS", 12, True, False, RGBColor(40, 40, 40)),
])

# SLIDE 2 — Welcome + Host
s = new_slide()
center_box(s, [
    label_line("Tonight's Host", 13),
    spacer(8),
    ("Terry Dwobeng", 60, True, False, WHITE),
    spacer(6),
    ("Founder, Property Labs", 24, False, False, WHITE_M),
    spacer(8),
    ("14 Units · £50k/month Revenue · Zero Own Capital", 18, False, False, RGBColor(120, 120, 120)),
])

# SLIDE 3 — Tonight is different
s = new_slide()
center_box(s, [
    ("Tonight is ", 52, True, False, WHITE),
    ("unlike anything", 52, True, False, GREEN),
    ("you've ever sat through.", 52, True, False, WHITE),
])

# SLIDE 4 — No slides. No pitch.
s = new_slide()
center_box(s, [
    ("No filler.", 40, False, False, WHITE_M),
    spacer(6),
    ("Just the exact system\nthat changed my life.", 52, True, False, WHITE),
])

# SLIDE 5 — Let Me Read Your Mind
s = new_slide()
center_box(s, [
    ("What if I told you", 20, True, False, GREEN),
    spacer(10),
    ("Let Me Read Your Mind", 64, True, False, WHITE),
])

# SLIDE 6 — Prediction 1a
s = new_slide()
left_box(s, [white_line("You have a good job.", 44)])

# SLIDE 7 — Prediction 1b
s = new_slide()
left_box(s, [
    white_line("You have a good job.", 44),
    spacer(8),
    green_line("But there's a voice that says…", 36),
    spacer(6),
    white_line('"Is this it?"', 52),
])

# SLIDE 8 — Prediction 2a
s = new_slide()
left_box(s, [white_line("You've thought about property.", 44)])

# SLIDE 9 — Prediction 2b
s = new_slide()
left_box(s, [
    white_line("You've thought about property.", 44),
    spacer(8),
    green_line("But every time you look into it…", 36),
    spacer(6),
    muted_line("You close the tab.", 36),
])

# SLIDE 10 — Prediction 3a
s = new_slide()
left_box(s, [white_line("Money wasn't handed to you.", 44)])

# SLIDE 11 — Prediction 3b
s = new_slide()
left_box(s, [
    white_line("Money wasn't handed to you.", 44),
    spacer(8),
    green_line("No trust fund.", 36),
    green_line("No Bank of Mum and Dad.", 36),
    spacer(6),
    muted_line("Just you.", 36),
])

# SLIDE 12 — Prediction 4a
s = new_slide()
left_box(s, [white_line("You're tired.", 52)])

# SLIDE 13 — Prediction 4b
s = new_slide()
left_box(s, [
    white_line("You're tired.", 52),
    spacer(8),
    green_line("Not physically.", 36),
    spacer(4),
    muted_line("Mentally.", 36),
])

# SLIDE 14 — Prediction 5a
s = new_slide()
left_box(s, [white_line("Deep down, you believe you're capable of more.", 40)])

# SLIDE 15 — Prediction 5b
s = new_slide()
left_box(s, [
    white_line("Deep down, you believe you're capable of more.", 40),
    spacer(8),
    green_line("You just don't have the system.", 40),
])

# SLIDE 16 — 57 nights on an air mattress
s = new_slide()
center_box(s, [
    muted_line("I know, because I was there.", 24),
    spacer(10),
    ("57 nights on an ", 56, True, False, WHITE),
    ("air mattress.", 56, True, False, GREEN),
])

# SLIDE 17 — Today?
s = new_slide()
center_box(s, [
    label_line("Today", 13),
    spacer(8),
    ("14 units.  £50,000/month revenue.", 48, True, False, WHITE),
    spacer(8),
    ("Zero of my own money.", 40, True, False, GREEN),
])

# SLIDE 18 — Zero.
big_number_slide("Zero.", "Not a penny of my own capital.")

# SLIDE 19 — The 3C System tease
s = new_slide()
center_box(s, [
    label_line("Tonight I'm revealing", 13),
    spacer(8),
    ("The 3C System", 64, True, False, GREEN),
    spacer(8),
    ("The exact framework I used.\nThe exact framework complete beginners are using right now.", 24, False, False, WHITE_M),
])

# SLIDE 20 — Complete beginners
s = new_slide()
center_box(s, [
    ("Complete beginners.", 56, True, False, GREEN),
    spacer(8),
    ("No experience. No capital. No connections.", 32, False, False, WHITE_M),
    spacer(6),
    ("All profitable inside 90 days.", 36, True, False, WHITE),
])

# ── ACT 2 · PAIN ─────────────────────────────────────────────────────────────

# SLIDE 21 — Act 2 intro
s = new_slide()
center_box(s, [
    ("But first.", 60, True, False, WHITE),
    spacer(10),
    ("Let's talk about the path\nyou were ", 32, False, False, WHITE_M),
    ("sold.", 32, True, False, GREEN),
])

# SLIDE 22 — The Traditional Path header
s = new_slide()
center_box(s, [
    label_line("The Traditional Path to Property", 13),
    spacer(8),
    ("Save → Mortgage → 1 Property", 44, True, False, WHITE),
    spacer(6),
    ("Wait → Repeat", 36, False, False, WHITE_M),
])

# SLIDE 23 — Average house price
s = new_slide()
center_box(s, [
    label_line("The Problem", 13),
    spacer(8),
    ("Average house price", 36, False, False, WHITE),
    spacer(6),
    ("£290,000", 88, True, False, GREEN),
])

# SLIDE 24 — Average salary
s = new_slide()
center_box(s, [
    ("Average UK salary", 36, False, False, WHITE_M),
    spacer(6),
    ("£34,000", 88, True, False, GREEN),
])

# SLIDE 25 — 4-6 years saving
s = new_slide()
center_box(s, [
    ("4–6 years", 56, True, False, GREEN),
    (" of aggressive saving.", 40, True, False, WHITE),
    spacer(10),
    ("For one property. One.", 24, False, False, WHITE_M),
])

# SLIDE 26 — Treadmill
s = new_slide()
center_box(s, [
    italic_white("That's not a wealth strategy.", 40),
    spacer(12),
    ("That's a treadmill.", 56, True, False, GREEN),
])

# SLIDE 27 — 5 Misconceptions header
section_divider("What's actually holding you back", "5 Misconceptions",
                "Every single one is keeping you stuck. Let's destroy them.")

# SLIDE 28 — Misconception #1 header
misconception("#1", "You need your own money to start.", "FALSE.")

# SLIDE 29 — Misc #1 proof
s = new_slide()
left_box(s, [
    muted_line("Other people's capital exists.", 36),
    spacer(10),
    ("Your credibility ", 44, True, False, WHITE),
    ("IS", 44, True, False, GREEN),
    (" the capital.", 44, True, False, WHITE),
])

# SLIDE 30 — Misconception #2
misconception("#2", "You need years of experience.", "FALSE.",
              "Ihab — zero experience. £25k raised. 2 Airbnbs live. 5 weeks.")

# SLIDE 31 — Misconception #3
misconception("#3", "Airbnb is saturated.", "FALSE.",
              "93,000 landlords left the market last year. The space just got bigger.")

# SLIDE 32 — Misconception #4
misconception("#4", "You need to quit your job.", "FALSE.",
              "Andrew — Masters student. Junior accountant. 3 units running. 5–10 hrs/week.")

# SLIDE 33 — Misconception #5
misconception("#5", "Now isn't the right time.", "FALSE.",
              "This is the best window in years. I'll prove it in Secret 3.")

# SLIDE 34 — Everything you've been told
s = new_slide()
center_box(s, [
    ("What if everything\nyou've been told…", 52, True, False, WHITE),
    spacer(10),
    ("is wrong?", 56, True, False, GREEN),
])

# SLIDE 35 — Different path
s = new_slide()
center_box(s, [
    ("What if there's a completely different path?", 36, False, False, WHITE_M),
    spacer(8),
    ("One that ", 44, True, False, WHITE),
    ("doesn't require", 44, True, False, GREEN),
    (" your own money.", 44, True, False, WHITE),
])

# SLIDE 36 — Let me show you
s = new_slide()
center_box(s, [
    ("Let me show you ", 52, True, False, WHITE),
    ("exactly", 52, True, False, GREEN),
    (" how I found it.", 52, True, False, WHITE),
])

# SLIDE 37 — Bridge to origin story
s = new_slide()
center_box(s, [
    label_line("The origin", 13),
    spacer(8),
    ("It started with ", 48, True, False, WHITE),
    ("an air mattress", 48, True, False, GREEN),
    (" and a question.", 48, True, False, WHITE),
])

# ── ACT 3 · POSITIONING — Origin Story ────────────────────────────────────────

# SLIDE 38 — South London
s = new_slide()
left_box(s, [
    muted_line("South London.", 44),
    spacer(8),
    white_line("Immigrant family.", 44),
    spacer(8),
    green_line("Nothing handed to us.", 44),
])

# SLIDE 39 — Did everything right
s = new_slide()
left_box(s, [
    white_line("Did everything right.", 44),
    spacer(10),
    muted_line("Kept my head down. Got the grades.", 32),
])

# SLIDE 40 — University of Warwick
s = new_slide()
center_box(s, [
    label_line("University of Warwick", 13),
    spacer(8),
    ("Law.", 80, True, False, GREEN),
])

# SLIDE 41 — Magic circle offer
s = new_slide()
center_box(s, [
    label_line("The offer", 13),
    spacer(8),
    ("Magic circle law firm.", 52, True, False, WHITE),
    spacer(6),
    ("Top 5 in the UK.", 28, False, False, WHITE_M),
])

# SLIDE 42 — Mum cried
s = new_slide()
center_box(s, [
    italic_white('"My mum cried."', 48),
    spacer(10),
    muted_line("She thought we'd finally made it.", 24),
])

# SLIDE 43 — Bigger cage
s = new_slide()
center_box(s, [
    muted_line("A bigger cage", 52),
    spacer(8),
    green_line("is still a cage.", 56),
])

# SLIDE 44 — 57 nights
big_number_slide("57", "nights on an air mattress.")

# SLIDE 45 — Staring at the ceiling
s = new_slide()
left_box(s, [
    white_line("Lying on that deflated mattress.", 40),
    spacer(8),
    muted_line("Staring at the ceiling.", 40),
    spacer(8),
    green_line("Knowing something had to change.", 40),
])

# SLIDE 46 — Kiyosaki quote
s = new_slide()
center_box(s, [
    italic_white('"The richest people don\'t work for money.\nMoney works for them."', 40),
    spacer(12),
    label_line("Robert Kiyosaki", 13),
])

# SLIDE 47 — But how?
s = new_slide()
center_box(s, [
    ("How do you get money\nto work for you…", 44, True, False, WHITE),
    spacer(8),
    green_line("when you don't have any?", 44),
])

# SLIDE 48 — That question changed everything
s = new_slide()
center_box(s, [
    muted_line("That question", 48),
    spacer(8),
    ("changed my ", 52, True, False, WHITE),
    ("life.", 52, True, False, GREEN),
])

# SLIDE 49 — Everyone said save
s = new_slide()
left_box(s, [
    muted_line('Everyone said:', 40),
    spacer(8),
    white_line('"Save up. Get a mortgage."', 44),
])

# SLIDE 50 — Law degree + air mattress
s = new_slide()
left_box(s, [
    white_line("I had a law degree.", 44),
    spacer(8),
    green_line("And an air mattress.", 44),
    spacer(8),
    muted_line("No savings. No contacts in property. Nothing.", 28),
])

# SLIDE 51 — What if my credibility IS the capital?
s = new_slide()
center_box(s, [
    italic_white("Then I asked a different question.", 32),
    spacer(12),
    ("What if my ", 48, True, False, WHITE),
    ("credibility", 48, True, False, GREEN),
    (" IS the capital?", 48, True, False, WHITE),
])

# SLIDE 52 — 3C System reveal
s = new_slide()
center_box(s, [
    label_line("The answer became", 13),
    spacer(10),
    ("C  →  C  →  C", 72, True, False, GREEN),
    spacer(8),
    ("Credibility  →  Capital  →  Cash Flow", 24, False, False, WHITE_M),
    spacer(10),
    ("The 3C System", 36, True, False, GREEN),
])

# SLIDE 53 — First property
s = new_slide()
center_box(s, [
    label_line("The result", 13),
    spacer(8),
    ("First property — profitable ", 44, True, False, WHITE),
    ("month one.", 44, True, False, GREEN),
])

# SLIDE 54 — Where I am now
s = new_slide()
center_box(s, [
    label_line("Today", 13),
    spacer(8),
    ("14 units.", 56, True, False, GREEN),
    spacer(6),
    ("£50k/month.  £16k profit.", 40, True, False, WHITE),
    spacer(6),
    ("Zero own money.", 36, True, False, GREEN),
])

# SLIDE 55 — Not just my story
s = new_slide()
center_box(s, [
    muted_line("This isn't just my story.", 40),
    spacer(8),
    ("This is a ", 52, True, False, WHITE),
    ("system.", 52, True, False, GREEN),
    spacer(6),
    muted_line("A system anyone can run.", 28),
])

# SLIDE 56 — Three secrets intro
s = new_slide()
center_box(s, [
    label_line("Tonight I'm sharing", 13),
    spacer(8),
    ("Three ", 48, True, False, WHITE),
    ("secrets", 48, True, False, GREEN),
    (" most people will never know.", 48, True, False, WHITE),
    spacer(8),
    muted_line("Let's start with Secret 1.", 24),
])

# ── ACT 4 · SECRET 1 — CREDIBILITY ────────────────────────────────────────────

# SLIDE 57 — SECRET 1 divider
section_divider("The first secret", "SECRET 1", "Your Credibility IS Your Capital")

# SLIDE 58 — Credibility as capital
s = new_slide()
center_box(s, [
    ("Your ", 52, True, False, WHITE),
    ("Credibility", 52, True, False, GREEN),
    (" IS Your Capital", 52, True, False, WHITE),
])

# SLIDE 59 — Private capital raising
s = new_slide()
left_box(s, [
    label_line("How the wealthy have always done it", 13),
    spacer(8),
    ("Private capital raising.", 52, True, False, WHITE),
    spacer(8),
    muted_line("The mechanism is not new.\nAccess to it, for people like us — that's new.", 24),
])

# SLIDE 60 — Credibility is not about being rich
s = new_slide()
center_box(s, [
    muted_line("Credibility is not about being rich.", 36),
    spacer(8),
    ("It's about ", 44, True, False, WHITE),
    ("packaging", 44, True, False, GREEN),
    (" what you already have.", 44, True, False, WHITE),
])

# SLIDE 61 — Do you have a degree?
s = new_slide()
left_box(s, [
    label_line("Ask yourself", 13),
    spacer(8),
    white_line("Do you have a degree?", 36),
    white_line("A professional job?", 36),
    white_line("A track record of any kind?", 36),
    spacer(8),
    green_line("You already have credibility.", 40),
])

# SLIDE 62 — Haven't packaged it
s = new_slide()
center_box(s, [
    muted_line("You already have it.", 48),
    spacer(8),
    ("You just haven't ", 48, True, False, WHITE),
    ("packaged it.", 48, True, False, GREEN),
])

# SLIDE 63 — The Credibility Pack
s = new_slide()
left_box(s, [
    label_line("The tool", 13),
    spacer(8),
    ("The Credibility Pack", 44, True, False, GREEN),
    spacer(10),
    ("→  Who you are", 28, False, False, WHITE),
    ("→  The strategy", 28, False, False, WHITE),
    ("→  The numbers", 28, False, False, WHITE),
    ("→  Legal structure", 28, False, False, WHITE),
    ("→  Risk mitigation", 28, False, False, WHITE),
])

# SLIDE 64 — Your network trusts you
s = new_slide()
center_box(s, [
    muted_line("Your network already trusts you.", 36),
    spacer(8),
    ("The Credibility Pack gives them\nthe ", 40, True, False, WHITE),
    ("rational justification.", 40, True, False, GREEN),
])

# SLIDE 65 — Case Study: IHAB header
case_study_header("Ihab", "Zero property experience.")

# SLIDE 66 — Ihab Week 1-2
s = new_slide()
left_box(s, [
    label_line("Ihab — Week 1–2", 13),
    spacer(8),
    white_line("Built his Credibility Pack.", 36),
    spacer(6),
    muted_line("Document. Strategy. Numbers. Legal framework.", 22),
])

# SLIDE 67 — Ihab Week 3
s = new_slide()
left_box(s, [
    label_line("Ihab — Week 3", 13),
    spacer(8),
    ("Approached ", 36, True, False, WHITE),
    ("3 people", 36, True, False, GREEN),
    (" in his network.", 36, True, False, WHITE),
    spacer(6),
    muted_line("Not cold contacts. People who already knew him.", 22),
])

# SLIDE 68 — Ihab Week 5
s = new_slide()
left_box(s, [
    label_line("Ihab — Week 5", 13),
    spacer(8),
    ("£25,000 committed.", 44, True, False, GREEN),
    spacer(6),
    muted_line("From people who trusted him. Not strangers.", 22),
])

# SLIDE 69 — Ihab Day 45
s = new_slide()
left_box(s, [
    label_line("Ihab — Day 45", 13),
    spacer(8),
    ("2 Airbnbs ", 36, True, False, WHITE),
    ("live.", 36, True, False, GREEN),
    spacer(6),
    ("Cash-flow positive.", 36, True, False, GREEN),
    spacer(6),
    muted_line("Zero property experience. Six weeks prior.", 22),
])

# SLIDE 70 — Two models
s = new_slide()
center_box(s, [
    label_line("Two ways to structure investor deals", 13),
    spacer(8),
    ("Equity vs Debt", 52, True, False, GREEN),
])

# SLIDE 71 — Equity model
s = new_slide()
left_box(s, [
    label_line("Model 1 — Equity", 13),
    spacer(8),
    ("Investor funds 100%.", 44, True, False, GREEN),
    spacer(6),
    white_line("You split profits.", 44),
    spacer(8),
    muted_line("You bring the system.\nThey bring the capital. Both win.", 22),
])

# SLIDE 72 — Debt model
s = new_slide()
left_box(s, [
    label_line("Model 2 — Debt", 13),
    spacer(8),
    ("Loan at agreed interest.", 44, True, False, GREEN),
    spacer(6),
    white_line("Pay back with returns.", 44),
    spacer(8),
    muted_line("Clean. Simple. Legally structured.\nInvestor gets a fixed return. You keep the rest.", 22),
])

# SLIDE 73 — Legal structures
s = new_slide()
center_box(s, [
    muted_line("Both models are", 40),
    spacer(6),
    ("legally ", 48, True, False, WHITE),
    ("structured.", 48, True, False, GREEN),
    spacer(8),
    muted_line("JV agreements. Loan agreements. Proper documentation.", 22),
])

# SLIDE 74 — Secret 1 recap
s = new_slide()
center_box(s, [
    label_line("Secret 1 Recap", 13),
    spacer(8),
    white_line("Package your credibility.", 44),
    spacer(6),
    green_line("Approach your network.", 40),
    spacer(6),
    ("Raise the capital.", 40, True, False, GREEN),
])

# SLIDE 75 — Transition to Secret 2
s = new_slide()
center_box(s, [
    muted_line("Now you have the capital.", 36),
    spacer(8),
    ("Secret 2 tells you ", 44, True, False, WHITE),
    ("exactly what to do with it.", 44, True, False, GREEN),
])

# ── ACT 5 · SECRET 2 — CAPITAL ────────────────────────────────────────────────

# SLIDE 76 — SECRET 2 divider
section_divider("The second secret", "SECRET 2",
                "You Don't Need Experience. You Need a System.")

# SLIDE 77 — You don't need experience
s = new_slide()
center_box(s, [
    muted_line("You don't need experience.", 48),
    spacer(8),
    ("You need a ", 52, True, False, WHITE),
    ("system.", 52, True, False, GREEN),
])

# SLIDE 78 — Capital Deployment Framework
s = new_slide()
center_box(s, [
    label_line("Introducing", 13),
    spacer(8),
    ("The Capital Deployment Framework", 44, True, False, GREEN),
    spacer(8),
    muted_line("Five steps. Rinse and repeat.", 24),
])

# SLIDE 79 — Step 1
s = new_slide()
left_box(s, [
    label_line("Step 1", 13),
    spacer(8),
    white_line("Deal Sourcing Team", 48),
    spacer(8),
    muted_line("They find the property.\nYou don't spend weekends on Rightmove.", 24),
])

# SLIDE 80 — Step 2
s = new_slide()
left_box(s, [
    label_line("Step 2", 13),
    spacer(8),
    white_line("Numbers Analysis", 48),
    spacer(8),
    muted_line("If it doesn't work on paper, we don't touch it.\nSimple rule. Protects everyone.", 24),
])

# SLIDE 81 — Step 3
s = new_slide()
left_box(s, [
    label_line("Step 3", 13),
    spacer(8),
    white_line("Property Setup", 48),
    spacer(8),
    muted_line("Furnishing. Photography. Optimised listing.\nPriced to win from day one.", 24),
])

# SLIDE 82 — Step 4
s = new_slide()
left_box(s, [
    label_line("Step 4", 13),
    spacer(8),
    white_line("Operations", 48),
    spacer(8),
    muted_line("Cleaning team. Maintenance. Automated guest messaging.\nYou're not the cleaner.", 24),
])

# SLIDE 83 — Step 5
s = new_slide()
left_box(s, [
    label_line("Step 5", 13),
    spacer(8),
    white_line("Optimisation", 48),
    spacer(8),
    muted_line("Dynamic pricing. Review strategy. Corporate contracts.\nRevenue grows over time.", 24),
])

# SLIDE 84 — You provide capital, system does rest
s = new_slide()
center_box(s, [
    muted_line("You provide the capital.", 36),
    spacer(8),
    ("The ", 52, True, False, WHITE),
    ("system", 52, True, False, GREEN),
    (" does the rest.", 52, True, False, WHITE),
])

# SLIDE 85 — Warm Outreach Framework
s = new_slide()
center_box(s, [
    label_line("A framework within the framework", 13),
    spacer(8),
    ("The Warm Outreach Framework", 44, True, False, GREEN),
    spacer(8),
    muted_line("The exact words to say when\napproaching investors in your network.", 24),
])

# SLIDE 86 — The script
s = new_slide()
center_box(s, [
    italic_white('"I\'m building an Airbnb portfolio\nand looking for investment partners."', 36),
    spacer(12),
    italic_white('"The investor puts in capital. I do all the work.\nWe split profits."', 28),
])

# SLIDE 87 — Case Study: ANDREW header
case_study_header("Andrew", "Masters student. Junior accountant. No savings.")

# SLIDE 88 — Andrew Month 1
s = new_slide()
left_box(s, [
    label_line("Andrew — Month 1", 13),
    spacer(8),
    ("Raised ", 36, True, False, WHITE),
    ("£13,000.", 40, True, False, GREEN),
    spacer(6),
    muted_line("Friends and former colleagues.\nWarm network. Credibility Pack.", 22),
])

# SLIDE 89 — Andrew Month 3
s = new_slide()
left_box(s, [
    label_line("Andrew — Month 3", 13),
    spacer(8),
    ("3 units live.", 40, True, False, GREEN),
    spacer(6),
    white_line("Cash flowing.", 36),
    spacer(6),
    muted_line("Still completing his Masters.\nStill working as an accountant.", 22),
])

# SLIDE 90 — Andrew 5-10 hours
big_number_slide("5–10", "hours per week. That's it.")

# SLIDE 91 — The system is the moat
s = new_slide()
center_box(s, [
    muted_line("Most people never launch because\nthey're trying", 32),
    spacer(8),
    ("to do it ", 52, True, False, WHITE),
    ("alone.", 52, True, False, GREEN),
])

# SLIDE 92 — You're not alone
s = new_slide()
center_box(s, [
    white_line("With Property Labs,", 44),
    spacer(8),
    green_line("you never are.", 48),
])

# SLIDE 93 — Secret 2 recap
s = new_slide()
center_box(s, [
    label_line("Secret 2 Recap", 13),
    spacer(8),
    white_line("A proven 5-step system\nreplaces years of experience.", 40),
    spacer(8),
    green_line("Anyone can run it.", 40),
])

# SLIDE 94 — Bridge to Secret 3
s = new_slide()
center_box(s, [
    muted_line("Credibility raised the capital.", 36),
    spacer(6),
    white_line("The system deployed it.", 36),
    spacer(6),
    green_line("Now let's talk about the cash flow.", 40),
])

# ── ACT 6 · SECRET 3 — CASH FLOW ─────────────────────────────────────────────

# SLIDE 95 — SECRET 3 divider
section_divider("The third secret", "SECRET 3",
                "The Biggest Opportunity Window in UK Airbnb History")

# SLIDE 96 — Biggest opportunity
s = new_slide()
center_box(s, [
    white_line("The Biggest Opportunity Window", 48),
    spacer(8),
    ("in UK Airbnb History.", 48, True, False, GREEN),
])

# SLIDE 97 — 93,000
big_number_slide("93,000", "UK landlords left the market last year.")

# SLIDE 98 — Largest exodus
s = new_slide()
center_box(s, [
    muted_line("The largest exodus", 48),
    spacer(8),
    ("in modern UK ", 48, True, False, WHITE),
    ("property history.", 48, True, False, GREEN),
])

# SLIDE 99 — What that means
s = new_slide()
center_box(s, [
    green_line("Less competition.", 44),
    spacer(6),
    ("More available properties.", 44, True, False, GREEN),
    spacer(6),
    white_line("Bigger margins for professionals.", 44),
])

# SLIDE 100 — Interest rates dropping
s = new_slide()
left_box(s, [
    label_line("Factor 2", 13),
    spacer(8),
    white_line("Interest Rates Are Dropping", 48),
    spacer(8),
    muted_line("Cheaper capital. Lower borrowing costs.\nBetter net returns on every unit.", 24),
])

# SLIDE 101 — New registration scheme
s = new_slide()
left_box(s, [
    label_line("Factor 3", 13),
    spacer(8),
    white_line("New Registration Scheme", 48),
    spacer(8),
    muted_line("The government is clearing out amateur hosts.\nCompliance requirements are rising. Professionals win.", 24),
])

# SLIDE 102 — Professional operators win
s = new_slide()
center_box(s, [
    muted_line("When the rules tighten,", 36),
    spacer(8),
    ("professional operators ", 48, True, False, WHITE),
    ("win.", 48, True, False, GREEN),
])

# SLIDE 103 — The equation
s = new_slide()
center_box(s, [
    label_line("The equation", 13),
    spacer(8),
    green_line("Less Competition", 32),
    white_line("+  Cheaper Capital", 32),
    green_line("+  Higher Barriers to Entry", 32),
    spacer(8),
    ("=  ", 44, True, False, WHITE),
    ("Bigger Moat", 44, True, False, GREEN),
])

# SLIDE 104 — The window
s = new_slide()
center_box(s, [
    muted_line("The window is open.", 44),
    spacer(8),
    green_line("It won't stay open forever.", 48),
])

# SLIDE 105 — Case Study: DYLAN header
case_study_header("Dylan", "Six-figure events company.\nWanted assets, not just income.")

# SLIDE 106 — Dylan 3 months
s = new_slide()
left_box(s, [
    label_line("Dylan — 3 months in", 13),
    spacer(8),
    ("5 Airbnbs.", 44, True, False, GREEN),
    spacer(6),
    white_line("All systemised. All running.", 36),
])

# SLIDE 107 — Dylan best month
s = new_slide()
center_box(s, [
    label_line("Dylan's best month", 13),
    spacer(8),
    ("£12,000", 100, True, False, GREEN),
])

# SLIDE 108 — Dylan Richard Mille
s = new_slide()
center_box(s, [
    muted_line("He bought a Richard Mille.", 36),
    spacer(8),
    ("Because the ", 48, True, False, WHITE),
    ("system works.", 48, True, False, GREEN),
])

# SLIDE 109 — Case Study: ELITE KLEANS header
case_study_header("Elite Kleans", "Cleaning company.\nWanted to own, not just service.")

# SLIDE 110 — Elite Kleans result
s = new_slide()
left_box(s, [
    label_line("Elite Kleans — Result", 13),
    spacer(8),
    white_line("1 unit sourced.", 36),
    spacer(6),
    ("£2,200/month.", 44, True, False, GREEN),
    spacer(6),
    muted_line("Service business became\nan asset-backed portfolio.", 22),
])

# SLIDE 111 — 3C Recap
s = new_slide()
center_box(s, [
    label_line("The full system", 13),
    spacer(10),
    ("C  →  C  →  C", 72, True, False, GREEN),
    spacer(8),
    ("Credibility  →  Capital  →  Cash Flow", 24, False, False, WHITE_M),
    spacer(10),
    ("This is how complete beginners did it.\nThis is how ", 24, False, False, WHITE_M),
    ("you", 24, True, False, GREEN),
    (" can do it.", 24, False, False, WHITE_M),
])

# SLIDE 112 — Capital Follows Conviction
s = new_slide()
center_box(s, [
    ("CAPITAL FOLLOWS CONVICTION", 72, True, False, GREEN),
])

# SLIDE 113 — What you need
s = new_slide()
center_box(s, [
    muted_line("You don't need money.", 40),
    spacer(8),
    ("You need ", 44, True, False, WHITE),
    ("credibility,", 44, True, False, GREEN),
    (" a system, and conviction.", 44, True, False, WHITE),
])

# ── ACT 7 · BRIDGE ────────────────────────────────────────────────────────────

# SLIDE 114 — The fork in the road
s = new_slide()
center_box(s, [
    label_line("Right now, in this moment", 13),
    spacer(8),
    ("You're standing at a ", 48, True, False, WHITE),
    ("fork in the road.", 48, True, False, GREEN),
])

# SLIDE 115 — Two paths
s = new_slide()
left_box(s, [
    label_line("Two paths forward", 13),
    spacer(8),
    ("PATH ONE", 20, True, False, GREEN),
    ("Leave tonight feeling inspired.", 26, False, False, WHITE_M),
    ("Feel great for 24 hours.", 26, False, False, WHITE_M),
    ("Back to normal by Monday.", 26, False, False, WHITE_M),
    spacer(10),
    ("PATH TWO", 20, True, False, GREEN),
    ("Decide tonight.", 26, False, False, WHITE),
    ("Take the next step.", 26, False, False, WHITE),
    ("12 months from now — different life.", 26, False, False, WHITE),
])

# SLIDE 116 — Path One outcome
s = new_slide()
left_box(s, [
    label_line("12 months on Path One", 13),
    spacer(8),
    white_line("Same salary.", 40),
    spacer(6),
    green_line("Same frustration.", 40),
    spacer(6),
    muted_line("Same scroll.", 40),
])

# SLIDE 117 — Path Two: Decide tonight
s = new_slide()
center_box(s, [
    white_line("Path Two starts with one decision.", 44),
    spacer(8),
    green_line("Tonight.", 64),
])

# SLIDE 118 — Take the next step
s = new_slide()
center_box(s, [
    muted_line("Not a massive leap.", 40),
    spacer(8),
    ("Just the ", 52, True, False, WHITE),
    ("next step.", 52, True, False, GREEN),
    spacer(8),
    muted_line("A conversation. 30 minutes. That's it.", 24),
])

# SLIDE 119 — Complete beginners have proven it
s = new_slide()
center_box(s, [
    ("Complete beginners", 40, True, False, GREEN),
    (" have already proven this works.", 40, True, False, WHITE),
    spacer(8),
    muted_line("The only missing piece", 40),
    spacer(6),
    green_line("is you.", 52),
])

# ── ACT 8 · OFFER STACK ───────────────────────────────────────────────────────

# SLIDE 120 — What's inside Property Labs
section_divider("Everything you need to launch", "Property Labs",
                "Let me show you exactly what's inside.")

# SLIDE 121 — Component 1
s = new_slide()
left_box(s, [
    label_line("Component 1", 13),
    spacer(8),
    white_line("7-Module Airbnb\nFoundations Course", 44),
    spacer(8),
    muted_line("Everything from deal structure to guest experience.\nThe full operating playbook.", 24),
])

# SLIDE 122 — Component 2
s = new_slide()
left_box(s, [
    label_line("Component 2", 13),
    spacer(8),
    white_line("Raising Finance Course", 44),
    spacer(6),
    green_line('"Nobody else teaches this."', 28),
    spacer(6),
    muted_line("How to approach investors. The Credibility Pack.\nLegal frameworks. Word-for-word scripts.", 22),
])

# SLIDE 123 — Component 3
s = new_slide()
left_box(s, [
    label_line("Component 3", 13),
    spacer(8),
    white_line("Legal Contract Pack", 44),
    spacer(8),
    muted_line("JV agreements. Rent-to-rent contracts.\nInvestor loan agreements. Professionally drafted. Ready to use.", 22),
])

# SLIDE 124 — Component 4
s = new_slide()
left_box(s, [
    label_line("Component 4", 13),
    spacer(8),
    white_line("Deal Sourcing Team Access", 44),
    spacer(6),
    green_line('"This is what makes us different."', 28),
    spacer(6),
    muted_line("A team that finds the properties for you.\nYou don't spend months searching Rightmove.", 22),
])

# SLIDE 125 — Component 5
s = new_slide()
left_box(s, [
    label_line("Component 5", 13),
    spacer(8),
    white_line("Cleaning & Maintenance Team", 44),
    spacer(8),
    muted_line("Your operational backbone.\nVetted. Reliable. Built into the system from day one.", 22),
])

# SLIDE 126 — Component 6
s = new_slide()
left_box(s, [
    label_line("Component 6", 13),
    spacer(8),
    white_line("Weekly Group Coaching Calls", 44),
    spacer(8),
    muted_line("Live sessions. Real questions. Real answers.\nPlus a community of people building alongside you.", 22),
])

# SLIDE 127 — Component 7
s = new_slide()
left_box(s, [
    label_line("Component 7", 13),
    spacer(8),
    white_line("The Full Toolkit", 44),
    spacer(8),
    muted_line("Investor outreach scripts. Pricing calculators.\nDeal analysis sheets. Listing templates.\nEverything pre-built.", 22),
])

# SLIDE 128 — Value Stack
s = new_slide()
left_box(s, [
    label_line("Everything inside Property Labs", 13),
    spacer(8),
    ("7-Module Foundations Course", 20, False, False, WHITE),
    ("£2,000", 20, True, False, GREEN),
    ("Raising Finance Course", 20, False, False, WHITE),
    ("£3,000", 20, True, False, GREEN),
    ("Legal Contract Pack", 20, False, False, WHITE),
    ("£1,500", 20, True, False, GREEN),
    ("Deal Sourcing Team Access", 20, False, False, WHITE),
    ("£5,000+", 20, True, False, GREEN),
    ("Cleaning & Maintenance Team", 20, False, False, WHITE),
    ("£2,000", 20, True, False, GREEN),
    ("Weekly Group Coaching + Community", 20, False, False, WHITE),
    ("£3,000", 20, True, False, GREEN),
    ("Full Toolkit", 20, False, False, WHITE),
    ("£500", 20, True, False, GREEN),
])

# SLIDE 129 — Total value
s = new_slide()
center_box(s, [
    label_line("If you got this separately", 13),
    spacer(8),
    ("£17,000+", 88, True, False, WHITE_M),
])

# SLIDE 130 — 4 Tiers header
s = new_slide()
center_box(s, [
    label_line("Choose your path", 13),
    spacer(8),
    ("4 Investment Levels", 56, True, False, GREEN),
    spacer(8),
    muted_line("Each designed for where you are right now.", 24),
])

# SLIDE 131 — Tier 1
s = new_slide()
left_box(s, [
    label_line("Tier 1", 13),
    spacer(8),
    white_line("Foundations Course", 48),
    spacer(6),
    muted_line("Self-study. The full operating blueprint.", 24),
    spacer(8),
    ("£500", 64, True, False, GREEN),
])

# SLIDE 132 — Tier 2
s = new_slide()
left_box(s, [
    label_line("Tier 2", 13),
    spacer(8),
    white_line("Both Courses", 48),
    spacer(6),
    muted_line("The full blueprint including the Raising Finance Course.\nThe complete picture.", 24),
    spacer(8),
    ("£1,000", 64, True, False, GREEN),
])

# SLIDE 133 — Tier 3
s = new_slide()
left_box(s, [
    label_line("Tier 3 — Most Popular", 13),
    spacer(8),
    white_line("Group Coaching", 48),
    spacer(6),
    muted_line("Both courses + deal sourcing team + cleaning team\n+ weekly calls + community + full toolkit.", 24),
    spacer(8),
    ("£3,000", 64, True, False, GREEN),
])

# SLIDE 134 — Tier 4
s = new_slide()
left_box(s, [
    label_line("Tier 4", 13),
    spacer(8),
    white_line("1:1 Coaching", 48),
    spacer(6),
    muted_line("Everything in Tier 3, plus weekly personal sessions\nwith Terry. Maximum accountability. Maximum results.", 24),
    spacer(8),
    ("£5,000", 64, True, False, GREEN),
])

# SLIDE 135 — 4 Tiers overview
s = new_slide()
left_box(s, [
    label_line("All four levels at a glance", 13),
    spacer(8),
    ("Foundations Course         £500", 26, False, False, WHITE),
    ("Both Courses               £1,000", 26, False, False, WHITE),
    ("Group Coaching ★           £3,000", 26, True, False, GREEN),
    ("1:1 Coaching               £5,000", 26, False, False, WHITE),
])

# SLIDE 136 — The Guarantee header
s = new_slide()
center_box(s, [
    label_line("Our promise to you", 13),
    spacer(8),
    ("The ", 52, True, False, WHITE),
    ("90-Day Guarantee", 52, True, False, GREEN),
])

# SLIDE 137 — Guarantee detail
s = new_slide()
center_box(s, [
    italic_white("Your first profitable Airbnb in 90 days.", 44),
    spacer(10),
    muted_line("Or we keep working with you.", 32),
    spacer(6),
    ("Free.", 40, True, False, GREEN),
    (" Until you do.", 40, True, False, GREEN),
])

# SLIDE 138 — Why we can offer this
s = new_slide()
center_box(s, [
    muted_line("We can offer this because", 36),
    spacer(8),
    ("the ", 48, True, False, WHITE),
    ("system works.", 48, True, False, GREEN),
    spacer(8),
    muted_line("25 students have already proved it.\nMultiple times over.", 22),
])

# SLIDE 139 — Book your strategy call
s = new_slide()
center_box(s, [
    label_line("The next step", 13),
    spacer(8),
    white_line("Book Your Free", 48),
    green_line("30-Minute Strategy Call", 52),
    spacer(8),
    muted_line("We map out your plan. You decide if\nProperty Labs is right for you. Zero pressure.", 22),
])

# SLIDE 140 — QR/Link
s = new_slide()
center_box(s, [
    label_line("Scan or go to", 13),
    spacer(8),
    white_line("propertylabs.co.uk/call", 52),
    spacer(10),
    ("Book Free Strategy Call →", 32, True, False, GREEN),
    spacer(6),
    muted_line("Limited spots available. First come, first served.", 16),
])

# ── ACT 9 · OBJECTION HANDLING ────────────────────────────────────────────────

# SLIDE 141 — Objections intro
s = new_slide()
center_box(s, [
    muted_line("Before you decide, let me answer", 36),
    spacer(8),
    ("what you're ", 48, True, False, WHITE),
    ("really", 48, True, False, GREEN),
    (" thinking.", 48, True, False, WHITE),
])

# SLIDE 142 — Objection: No money
objection_slide('"I don\'t have the money to invest."',
                'Capital raising IS the system.\nThat\'s literally what we teach first.')

# SLIDE 143 — Objection: Will it work for me
objection_slide('"Will this actually work for someone like me?"',
                'Ihab. Andrew. Dylan. Elite Kleans.\n25 different people. Different backgrounds. Same system. Same results.')

# SLIDE 144 — Objection: Scared
objection_slide('"I\'m scared of making the wrong move."',
                'Educated risk with structured support is not reckless.\nStaying stuck forever is.')

# SLIDE 145 — Objection: I'll start later
objection_slide('"I\'ll think about it and maybe start later."',
                'Every month you wait = £2,000–£4,000 of income lost.\n6 months = £12,000 gone.')

# SLIDE 146 — Cost of waiting
big_number_slide("£12,000", "lost in 6 months of inaction.")

# SLIDE 147 — Objection: Airbnb saturated
objection_slide('"Isn\'t Airbnb saturated by now?"',
                '93,000 landlords left. Amateurs cleared out.\nCorporate demand at all-time high. Saturated? The opposite.')

# SLIDE 148 — Objection: Not experienced enough
objection_slide('"I\'m not experienced enough."',
                'I started on an air mattress. Ihab had zero property experience.\nExperience isn\'t the entry requirement.')

# SLIDE 149 — Objection: Need to think about it
objection_slide('"I just need a bit more time to think."',
                'Pain of discipline is temporary.\nPain of regret is permanent.\nThe people who acted last year are cash flowing today.')

# SLIDE 150 — Objection: Can't do it with my job
objection_slide('"I can\'t fit this around my job."',
                'Andrew. Masters student. Accountant. 3 units running.\n5–10 hours per week. The system is built around your life.')

# SLIDE 151 — One more question
s = new_slide()
center_box(s, [
    muted_line("Any other question you have", 40),
    spacer(8),
    ("ask it on the ", 44, True, False, WHITE),
    ("strategy call.", 44, True, False, GREEN),
])

# SLIDE 152 — CTA repeat
s = new_slide()
center_box(s, [
    label_line("Right now — while you still feel it", 13),
    spacer(8),
    green_line("Book Your Free Strategy Call", 48),
    spacer(8),
    white_line("propertylabs.co.uk/call", 44),
    spacer(8),
    ("Book Now →", 28, True, False, GREEN),
])

# ── ACT 10 · CLOSE ────────────────────────────────────────────────────────────

# SLIDE 153 — Two pictures intro
s = new_slide()
center_box(s, [
    muted_line("Before you decide,", 40),
    spacer(8),
    ("let me paint you ", 48, True, False, WHITE),
    ("two pictures.", 48, True, False, GREEN),
])

# SLIDE 154 — Picture One header
s = new_slide()
center_box(s, [
    label_line("Picture One", 13),
    spacer(8),
    muted_line("April 2027.", 64),
])

# SLIDE 155 — Picture One detail
s = new_slide()
left_box(s, [
    label_line("April 2027 — Path One", 13),
    spacer(8),
    muted_line("Same salary.", 40),
    muted_line("Same rent.", 40),
    muted_line("Same Sunday dread.", 40),
    spacer(6),
    green_line("Same scroll. Wondering what if.", 36),
])

# SLIDE 156 — Picture Two header
s = new_slide()
center_box(s, [
    label_line("Picture Two", 13),
    spacer(8),
    ("April 2027.", 64, True, False, GREEN),
])

# SLIDE 157 — Picture Two: Airbnb live
s = new_slide()
left_box(s, [
    label_line("April 2027 — Path Two", 13),
    spacer(8),
    white_line("First Airbnb live.", 40),
    spacer(6),
    ("£2,000–£4,000/month coming in.", 40, True, False, GREEN),
    spacer(6),
    white_line("Second unit loading.", 40),
])

# SLIDE 158 — Parents know
s = new_slide()
center_box(s, [
    muted_line("Your parents don't just think you've done well.", 32),
    spacer(8),
    ("They ", 52, True, False, WHITE),
    ("know.", 52, True, False, GREEN),
])

# SLIDE 159 — You feel free
s = new_slide()
center_box(s, [
    white_line("For the first time,", 44),
    spacer(8),
    green_line("you feel free.", 56),
])

# SLIDE 160 — Final 3C recap
s = new_slide()
center_box(s, [
    label_line("The system that makes it possible", 13),
    spacer(10),
    ("C  →  C  →  C", 72, True, False, GREEN),
    spacer(8),
    ("Credibility  →  Capital  →  Cash Flow", 24, False, False, WHITE_M),
])

# SLIDE 161 — Capital Follows Conviction (final)
s = new_slide()
center_box(s, [
    ("CAPITAL FOLLOWS\nCONVICTION", 80, True, False, GREEN),
])

# SLIDE 162 — The title
s = new_slide()
center_box(s, [
    muted_line("Tonight's title was:", 24),
    spacer(8),
    white_line('"The 3C System: How Complete Beginners Are\nLaunching Profitable Airbnbs in 90 Days…"', 32),
])

# SLIDE 163 — And How You Can Be Next
s = new_slide()
center_box(s, [
    white_line("I've shown you it works for complete beginners.", 32),
    spacer(6),
    white_line("I've shown you the system.", 32),
    spacer(8),
    green_line("The only missing piece is you.", 48),
])

# SLIDE 164 — Final CTA
s = new_slide()
center_box(s, [
    ("BOOK YOUR FREE STRATEGY CALL", 56, True, False, GREEN),
    spacer(8),
    white_line("propertylabs.co.uk/call", 44),
    spacer(8),
    ("Book Now — It's Free →", 28, True, False, GREEN),
])

# SLIDE 165 — One year from now
s = new_slide()
center_box(s, [
    italic_white('"One year from now, you\'ll wish\nyou\'d started tonight."', 48),
])

# ── FINAL: Q&A + End ─────────────────────────────────────────────────────────

# SLIDE 166 — Q&A
section_divider("Let's talk", "Q&A", "Your questions, answered. Fire away.")

# SLIDE 167 — While you think of questions
s = new_slide()
center_box(s, [
    muted_line("While you think of questions,", 36),
    spacer(8),
    ("remember — the strategy call is ", 44, True, False, WHITE),
    ("free.", 44, True, False, GREEN),
    spacer(8),
    muted_line("No commitment. No pressure. Just a plan.", 24),
])

# SLIDE 168 — The link again
s = new_slide()
center_box(s, [
    label_line("Book right now while we talk", 13),
    spacer(8),
    white_line("propertylabs.co.uk/call", 52),
    spacer(8),
    ("Secure Your Spot →", 32, True, False, GREEN),
])

# SLIDE 169 — Thank you
s = new_slide()
center_box(s, [
    label_line("Thank you for being here tonight", 13),
    spacer(8),
    white_line("Property Labs", 56),
    spacer(8),
    muted_line("Terry Dwobeng", 28),
    spacer(10),
    ("PROPERTY LABS · propertylabs.co.uk", 12, True, False, RGBColor(40, 40, 40)),
])

# SLIDE 170 — Final link slide
s = new_slide()
center_box(s, [
    green_line("Book Your Free Strategy Call", 52),
    spacer(8),
    white_line("propertylabs.co.uk/call", 44),
    spacer(10),
    muted_line("14 units. £50k/month. Zero own capital.\nThe system is proven. Your turn.", 18),
])

# ── SAVE ─────────────────────────────────────────────────────────────────────
out = "/Users/oliverknight/Jarvis/clients/Terry/webinar-slides/property-labs-webinar.pptx"
prs.save(out)

total = len(prs.slides)
print(f"Saved {total} slides → {out}")
